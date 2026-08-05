import os
import json
import logging
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt, Emu
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from docx2pdf import convert
from .base import BaseStage
from ..config import Config

logger = logging.getLogger(__name__)


# DOCX & PDF REPORT GENERATION
# ═══════════════════════════════════════════════════════


# ═══════════════════════════════════════════════════════
# DOCX REPORT (python-docx)
# ═══════════════════════════════════════════════════════
def generate_docx(state: dict, out_path: str):
    """Generate a comprehensive DOCX Teacher Guide from TKP state."""
    doc = Document()

    classification = state.get("classification") or {}
    knowledge = state.get("knowledge") or {}
    lesson_plan = state.get("lesson_plan") or {}
    period_contents = state.get("period_contents") or []
    activities = state.get("activities") or []
    assessments = state.get("ab_test_assessment") or {}
    gap_data = state.get("gap_analysis") or {}
    validation = state.get("validation") or {}

    subject = classification.get("subject", "General Curriculum")
    topic = classification.get("topic", "Teacher Knowledge Package")
    grade = classification.get("target_grade", classification.get("grade_level", "K-12"))
    board = classification.get("curriculum_board", "CBSE/NCERT")
    total_periods = lesson_plan.get("total_periods", len(period_contents) or 3)

    # ── Title ──
    title = doc.add_heading(f"{subject}: {topic}", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.LEFT
    subtitle = doc.add_paragraph(f"Grade {grade}  |  Board: {board}  |  {total_periods} Teaching Periods")
    subtitle.style.font.color.rgb = RGBColor(100, 116, 139)
    doc.add_paragraph("")

    # ── 1. Learning Objectives ──
    objs = knowledge.get("learning_objectives") or []
    if objs:
        doc.add_heading("1. Core Learning Objectives", level=1)
        for o in objs:
            text = o if isinstance(o, str) else o.get("objective", str(o))
            bloom = "" if isinstance(o, str) else o.get("blooms_level", "")
            p = doc.add_paragraph(text, style='List Bullet')
            if bloom:
                run = p.add_run(f"  [Bloom: {bloom}]")
                run.italic = True
                run.font.color.rgb = RGBColor(99, 102, 241)

    # ── 2. Key Concepts ──
    concepts = knowledge.get("concepts") or []
    definitions = knowledge.get("definitions") or []
    formulae = knowledge.get("formulae") or []
    if concepts or definitions:
        doc.add_heading("2. Key Concepts, Definitions & Formulae", level=1)
        if concepts:
            for c in concepts:
                name = c if isinstance(c, str) else c.get("name", "")
                desc = "" if isinstance(c, str) else c.get("description", "")
                p = doc.add_paragraph(style='List Bullet')
                run = p.add_run(f"{name}: ")
                run.bold = True
                p.add_run(str(desc))
        if definitions:
            for d in definitions:
                term = d if isinstance(d, str) else d.get("term", "")
                defn = "" if isinstance(d, str) else d.get("definition", "")
                p = doc.add_paragraph(style='List Bullet')
                run = p.add_run(f"{term}: ")
                run.bold = True
                p.add_run(str(defn))
        if formulae:
            doc.add_heading("Formulae", level=2)
            for f in formulae:
                name = f if isinstance(f, str) else f.get("name", "")
                expr = "" if isinstance(f, str) else (f.get("latex", "") or f.get("plain_text", ""))
                p = doc.add_paragraph(style='List Bullet')
                run = p.add_run(f"{name}:  ")
                run.bold = True
                run2 = p.add_run(str(expr))
                run2.font.name = "Courier New"

    # ── 3. Lesson Plan ──
    periods = lesson_plan.get("periods") or []
    if periods:
        doc.add_heading("3. Multi-Period Lesson Plan", level=1)
        # Summary table
        table = doc.add_table(rows=1, cols=4)
        table.style = 'Light Grid Accent 1'
        hdr = table.rows[0].cells
        hdr[0].text = "Period"
        hdr[1].text = "Title"
        hdr[2].text = "Objectives"
        hdr[3].text = "Methodology"
        for p in periods:
            row = table.add_row().cells
            row[0].text = str(p.get("period_number", 1))
            row[1].text = str(p.get("title", ""))
            row[2].text = ", ".join(str(o) for o in p.get("learning_objectives", []))
            row[3].text = str(p.get("teaching_methodology", ""))
        doc.add_paragraph("")

    # ── 4. Teacher Scripts ──
    if period_contents:
        doc.add_heading("4. Teacher Delivery Scripts", level=1)
        for pc in period_contents:
            p_num = pc.get("period_number", 1)
            doc.add_heading(f"Period {p_num}", level=2)

            entry = pc.get("entry_ticket", {})
            if entry and entry.get("question"):
                p = doc.add_paragraph()
                run = p.add_run("Entry Ticket: ")
                run.bold = True
                run.font.color.rgb = RGBColor(99, 102, 241)
                p.add_run(str(entry["question"]))

            script = pc.get("teacher_script", "")
            if script:
                doc.add_paragraph(str(script))

            bb = pc.get("blackboard_notes", "")
            if bb:
                p = doc.add_paragraph()
                run = p.add_run("Blackboard Notes: ")
                run.bold = True
                p.add_run(str(bb))

            exit_t = pc.get("exit_ticket", {})
            if exit_t and exit_t.get("question"):
                p = doc.add_paragraph()
                run = p.add_run("Exit Ticket: ")
                run.bold = True
                run.font.color.rgb = RGBColor(234, 88, 12)
                p.add_run(str(exit_t["question"]))

    # ── 5. Activities ──
    if activities:
        doc.add_heading("5. Classroom Activities", level=1)
        for act in activities:
            title_str = act.get("title", "Activity")
            doc.add_heading(f"{title_str} ({act.get('duration_minutes', 15)} mins)", level=2)
            if act.get("student_instructions"):
                p = doc.add_paragraph()
                run = p.add_run("Student Instructions: ")
                run.bold = True
                p.add_run(str(act["student_instructions"]))
            teacher_inst = act.get("teacher_instructions", [])
            if teacher_inst:
                p = doc.add_paragraph()
                run = p.add_run("Teacher Guidance: ")
                run.bold = True
                p.add_run("; ".join(str(t) for t in teacher_inst))
            materials = act.get("materials_needed", [])
            if materials:
                doc.add_paragraph(f"Materials: {', '.join(str(m) for m in materials)}")

    # ── 6. Assessments ──
    doc.add_heading("6. A/B Test Assessments", level=1)
    for var_key, var_label in [("variant_a", "Variant A — Standard"), ("variant_b", "Variant B — Deep Reasoning")]:
        variant = assessments.get(var_key, {})
        doc.add_heading(var_label, level=2)

        mcqs = variant.get("mcqs") or []
        if mcqs:
            doc.add_heading("MCQs", level=3)
            for i, q in enumerate(mcqs):
                p = doc.add_paragraph()
                run = p.add_run(f"Q{i+1}. {q.get('question', '')}")
                run.bold = True
                for oIdx, opt in enumerate(q.get("options", [])):
                    doc.add_paragraph(f"   ({chr(65+oIdx)}) {opt}")
                ans_p = doc.add_paragraph()
                ans_run = ans_p.add_run(f"Answer: {q.get('correct_option', '')} — {q.get('explanation', '')}")
                ans_run.font.color.rgb = RGBColor(16, 185, 129)

        short_ans = variant.get("short_answer") or []
        if short_ans:
            doc.add_heading("Short Answer", level=3)
            for i, q in enumerate(short_ans):
                p = doc.add_paragraph()
                run = p.add_run(f"Q{i+1}. {q.get('question', '')}")
                run.bold = True
                ans_p = doc.add_paragraph()
                ans_run = ans_p.add_run(f"Model Answer: {q.get('model_answer', '')}")
                ans_run.font.color.rgb = RGBColor(16, 185, 129)

    # ── 7. Gap Analysis ──
    gaps_list = gap_data.get("gaps") or []
    if gaps_list:
        doc.add_heading("7. Learning Gap Analysis & Remediation", level=1)
        table = doc.add_table(rows=1, cols=4)
        table.style = 'Light Grid Accent 1'
        hdr = table.rows[0].cells
        hdr[0].text = "Concept"
        hdr[1].text = "Misconception"
        hdr[2].text = "Diagnostic Question"
        hdr[3].text = "Remedial Action"
        for g in gaps_list:
            row = table.add_row().cells
            row[0].text = str(g.get("concept", ""))
            row[1].text = str(g.get("misconception", ""))
            row[2].text = str(g.get("diagnostic_question", ""))
            row[3].text = str(g.get("remedial_action", ""))

    # ── 8. Validation ──
    if validation:
        doc.add_heading("8. Quality Validation Report", level=1)
        completeness = validation.get("completeness_score", 0)
        consistency = validation.get("consistency_score", 0)
        overall = round((completeness + consistency) / 2) if completeness and consistency else 'N/A'
        doc.add_paragraph(f"Overall Score: {overall}/100")
        
        flags = validation.get('hallucination_flags', 0)
        if isinstance(flags, list):
            doc.add_paragraph(f"Hallucination Flags: {len(flags)}")
            for flag in flags:
                flag_text = flag.get("description", str(flag)) if isinstance(flag, dict) else str(flag)
                doc.add_paragraph(f"  - {flag_text}")
        elif isinstance(flags, dict):
            doc.add_paragraph(f"Hallucination Flags: 1")
            flag_text = flags.get("description", str(flags))
            doc.add_paragraph(f"  - {flag_text}")
        else:
            doc.add_paragraph(f"Hallucination Flags: {flags}")
            
        issues = validation.get("issues", [])
        for issue in issues:
            text = issue if isinstance(issue, str) else issue.get("description", str(issue))
            doc.add_paragraph(f"  - {text}")

    doc.save(out_path)


# ═══════════════════════════════════════════════════════
# PPTX REPORT (python-pptx)
# ═══════════════════════════════════════════════════════
def generate_pptx(state: dict, out_path: str):
    """Generate a PPTX presentation from TKP state."""
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    classification = state.get("classification") or {}
    knowledge = state.get("knowledge") or {}
    lesson_plan = state.get("lesson_plan") or {}
    period_contents = state.get("period_contents") or []
    activities = state.get("activities") or []
    assessments = state.get("ab_test_assessment") or {}
    gap_data = state.get("gap_analysis") or {}

    subject = classification.get("subject", "General Curriculum")
    topic = classification.get("topic", "Teacher Knowledge Package")
    grade = classification.get("target_grade", classification.get("grade_level", "K-12"))
    board = classification.get("curriculum_board", "CBSE/NCERT")
    total_periods = lesson_plan.get("total_periods", len(period_contents) or 3)

    def _create_single_slide(title_text, body_text):
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        if title:
            title.text = str(title_text)
        body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if body and body_text:
            body.text = str(body_text)
        return slide

    def add_slide(title_text, body_text=""):
        if body_text:
            # Strip multiple newlines so PPTX doesn't render empty bullets
            lines = [line.strip() for line in str(body_text).split("\n") if line.strip()]
            body_text = "\n".join(lines)
            
        chunk_size = 400 # Max characters before text overflows slide
        if not body_text or len(body_text) <= chunk_size:
            return _create_single_slide(title_text, body_text)

        # Chunk the text so it doesn't spill off the bottom
        lines = body_text.split('\n')
        current_chunk = ""
        part = 1
        for line in lines:
            if len(current_chunk) + len(line) > chunk_size and current_chunk:
                _create_single_slide(f"{title_text} (Part {part})", current_chunk.strip())
                current_chunk = line + "\n"
                part += 1
            else:
                current_chunk += line + "\n"
        if current_chunk.strip():
            _create_single_slide(f"{title_text} (Part {part})", current_chunk.strip())
        return None

    def add_title_slide(title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        if title:
            title.text = str(title_text)
        subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if subtitle:
            subtitle.text = str(subtitle_text)
        return slide

    # ── Slide 1: Title ──
    add_title_slide(
        f"{subject}: {topic}",
        f"Grade {grade}  |  Board: {board}  |  {total_periods} Teaching Periods\nTeacher Knowledge Package — AI Generated"
    )

    # ── Slide 2: Learning Objectives ──
    objs = knowledge.get("learning_objectives", [])
    if objs:
        obj_text = "\n".join(
            f"• {o if isinstance(o, str) else o.get('objective', str(o))}" for o in objs[:8]
        )
        add_slide("Core Learning Objectives", obj_text)

    # ── Slide 3: Key Concepts ──
    concepts = knowledge.get("concepts", [])
    if concepts:
        concept_text = "\n".join(
            f"• {c if isinstance(c, str) else c.get('name', '')}: {'' if isinstance(c, str) else c.get('description', '')}"
            for c in concepts[:8]
        )
        add_slide("Key Concepts & Definitions", concept_text)

    # ── Slides: Lesson Plan per period ──
    periods = lesson_plan.get("periods", [])
    for p in periods[:10]:
        p_num = p.get("period_number", 1)
        p_title = p.get("title", f"Period {p_num}")
        body = f"Objectives: {', '.join(str(o) for o in p.get('learning_objectives', []))}\n"
        body += f"Methodology: {p.get('teaching_methodology', '')}\n"
        body += f"Concepts: {', '.join(str(c) for c in p.get('concepts_covered', []))}"
        add_slide(f"Period {p_num}: {p_title}", body)

    # ── Slides: Teacher Scripts ──
    for pc in period_contents[:5]:
        p_num = pc.get("period_number", 1)
        script = str(pc.get("teacher_script", ""))
        add_slide(f"Period {p_num} — Teacher Script", script)

    # ── Slide: Activities ──
    if activities:
        act_text = "\n\n".join(
            f"• {act.get('title', 'Activity')} ({act.get('duration_minutes', 15)} mins): "
            f"{act.get('student_instructions', '')[:200]}"
            for act in activities[:6]
        )
        add_slide("Classroom Activities", act_text)

    # ── Slide: Assessment Variant A ──
    var_a = assessments.get("variant_a", {})
    mcqs_a = var_a.get("mcqs", [])
    if mcqs_a:
        mcq_text = "\n\n".join(
            f"Q{i+1}. {q.get('question', '')}\nAnswer: {q.get('correct_option', '')}"
            for i, q in enumerate(mcqs_a[:5])
        )
        add_slide("Assessment — Variant A (Standard)", mcq_text)

    # ── Slide: Assessment Variant B ──
    var_b = assessments.get("variant_b", {})
    mcqs_b = var_b.get("mcqs", [])
    if mcqs_b:
        mcq_text = "\n\n".join(
            f"Q{i+1}. {q.get('question', '')}\nAnswer: {q.get('correct_option', '')}"
            for i, q in enumerate(mcqs_b[:5])
        )
        add_slide("Assessment — Variant B (Deep Reasoning)", mcq_text)

    # ── Slide: Gap Analysis ──
    gaps_list = gap_data.get("gaps", [])
    if gaps_list:
        gap_text = "\n\n".join(
            f"• {g.get('concept', '')}: {g.get('misconception', '')}\n  Remediation: {g.get('remedial_action', '')}"
            for g in gaps_list[:5]
        )
        add_slide("Learning Gap Analysis", gap_text)

    # ── Final Slide ──
    add_title_slide("Thank You", "Generated by Teacher AI Platform\nPowered by a 10-Stage AI Pipeline")

    prs.save(out_path)


# ═══════════════════════════════════════════════════════
# PUBLISHING STAGE (orchestrates all three)
# ═══════════════════════════════════════════════════════
class PublishingStage(BaseStage):
    """
    Stage 10: Formatting, PDF/DOCX/PPTX Document Generation and Packaging.
    """
    def execute(self, state: dict) -> dict:
        job_id = self.document_id
        output_dir = Config.OUTPUT_FOLDER
        os.makedirs(output_dir, exist_ok=True)

        pdf_path = os.path.join(output_dir, f"TKP_{job_id}.pdf")
        docx_path = os.path.join(output_dir, f"TKP_{job_id}.docx")
        pptx_path = os.path.join(output_dir, f"TKP_{job_id}.pptx")
        json_path = os.path.join(output_dir, f"TKP_{job_id}.json")

        results = {"format": "JSON+PDF+DOCX+PPTX", "version": "1.0.0", "ready_for_export": True}

        # JSON master file
        try:
            with open(json_path, 'w', encoding='utf-8') as f:
                json.dump(state, f, indent=2, ensure_ascii=False, default=str)
            results["json_path"] = json_path
            logger.info(f"Generated JSON at: {json_path}")
        except Exception as e:
            logger.error(f"JSON export error: {e}")

        # 1. Generate DOCX first
        try:
            generate_docx(state, docx_path)
            results["docx_path"] = docx_path
            logger.info(f"Generated DOCX at: {docx_path}")
        except Exception as e:
            logger.error(f"DOCX export error: {e}")

        # 2. Generate PDF by converting DOCX using docx2pdf
        try:
            from docx2pdf import convert
            # Convert DOCX to PDF (requires MS Word on Windows/Mac)
            convert(docx_path, pdf_path)
            results["pdf_path"] = pdf_path
            logger.info(f"Generated PDF at: {pdf_path}")
        except Exception as e:
            logger.error(f"PDF export error (docx2pdf): {e}")

        # PPTX
        try:
            generate_pptx(state, pptx_path)
            results["pptx_path"] = pptx_path
            logger.info(f"Generated PPTX at: {pptx_path}")
        except Exception as e:
            logger.error(f"PPTX export error: {e}")

        return results
